import os
from pathlib import Path
from pptx import Presentation
import openai
from moviepy.editor import *
from pydub import AudioSegment

# Set your OPENAI_API_KEY in your environment variables here
# or otherwise assuming you've set the key outside of this script
#os.environ['OPENAI_API_KEY'] = ''

# Set path of the PPTX file and the folder containing the TIFF images for each slide
PPTX_PATH = 'p6.pptx'
IMAGE_FOLDER_PATH = './slides'  

# Set the gap (audio silence) between slides in milliseconds
GAP_BETWEEN_SLIDES_MS = 800
# Set the export location for the final video
MOV_EXPORT_PATH = "./output/presentation_video_export.mp4"
# Set the export FPS for the final video
MOV_EXPORT_FPS = 6

# TTS voice name
TTS_VOICE_NAME = "echo"

# Set the folder path to store the generated audio files
AUDIO_FOLDER_PATH = "./audio" 

client = openai.OpenAI()

def extract_speaker_notes(pptx_path):
    prs = Presentation(pptx_path)
    notes_texts = []
    for slide in prs.slides:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ''
        notes_texts.append(notes_text)
    return notes_texts


def text_to_speech(text, slide_number):
    speech_file_path = Path(AUDIO_FOLDER_PATH,f'slide_{slide_number}.mp3')
    response = client.audio.speech.create(
      model="tts-1",
      voice=TTS_VOICE_NAME,
      input=text
    )
    
    # Stream the TTS response to an audio file
    temp_audio_path = Path(f"temp_{slide_number}.mp3")
    response.stream_to_file(str(temp_audio_path))
    
    # Load the generated audio
    audio_clip = AudioSegment.from_file(str(temp_audio_path))
    
    # Create silence in ms to add at the end of the audio clip
    silence = AudioSegment.silent(duration=GAP_BETWEEN_SLIDES_MS)
    
    # Concatenate the audio clip with silence
    final_audio = audio_clip + silence
    
    # Export the concatenated audio
    final_audio.export(speech_file_path, format="mp3")
    
    # Cleanup the temporary file
    temp_audio_path.unlink()
    
    return speech_file_path


def create_video_presentation(image_folder, audio_paths):
    slides = sorted(Path(image_folder).glob('*.tiff'), key=os.path.getmtime)
    clips = []
    for slide_path, audio_path in zip(slides, audio_paths):
        audio_clip = AudioFileClip(str(audio_path))
        img_clip = ImageClip(str(slide_path)).set_duration(audio_clip.duration)
        img_clip = img_clip.set_audio(audio_clip)
        clips.append(img_clip)
    
    final_clip = concatenate_videoclips(clips, method="compose")
    final_clip.write_videofile(MOV_EXPORT_PATH, fps=MOV_EXPORT_FPS)

def main(pptx_path, image_folder):
    notes_texts = extract_speaker_notes(pptx_path)
    audio_paths = []
    for i, notes_text in enumerate(notes_texts, start=1):
        if notes_text:  # Proceed only if there is text
            audio_path = text_to_speech(notes_text, i)
            audio_paths.append(audio_path)
        else:
            # If no notes for a slide, consider handling silence or default audio
            audio_paths.append(None)  # Placeholder for slides without notes

    # Filter out slides without audio
    audio_paths = [path for path in audio_paths if path is not None]
    
    create_video_presentation(image_folder, audio_paths)


main(PPTX_PATH, IMAGE_FOLDER_PATH)